"""
报告导出工具 - 支持 Markdown / PDF / PPT 三种格式
"""
import json
from pathlib import Path
from datetime import datetime
from config import settings
from logger import setup_logger

logger = setup_logger("tools.report_exporter")

# 导出目录
EXPORT_DIR = settings.data_dir / "exports"


def _ensure_export_dir():
    EXPORT_DIR.mkdir(parents=True, exist_ok=True)


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


class ReportExporter:

    def export_markdown(self, report_data: dict) -> str:
        """导出 Markdown 格式报告，返回文件路径"""
        _ensure_export_dir()
        filename = f"decision_report_{_timestamp()}.md"
        filepath = EXPORT_DIR / filename
        logger.info(f"导出Markdown报告: {filename}")

        top10 = report_data.get("top10", [])
        core_nodes = report_data.get("coreNodes", [])
        not_recommended = report_data.get("notRecommended", [])
        decision_basis = report_data.get("decisionBasis", [])
        risk_warnings = report_data.get("riskWarnings", [])
        summary = report_data.get("summary", {})

        lines = [
            "# 需求优先级决策报告",
            f"\n> 生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}",
            "\n---",
            "\n## 一、总览",
            f"- 需求总数：{summary.get('totalRequirements', 0)}",
            f"- P0 需求：{summary.get('p0Count', 0)}",
            f"- 高风险需求：{summary.get('highRiskCount', 0)}",
            f"- 平均得分：{summary.get('avgScore', 0)}",
            "\n## 二、本轮优先级 Top 10",
            "| 排名 | 需求名称 | 优先级 | 总分 |",
            "| ---- | -------- | ------ | ---- |",
        ]

        for i, req in enumerate(top10, 1):
            lines.append(f"| {i} | {req.get('name', '')} | {req.get('priority', '')} | {req.get('totalScore', 0)} |")

        lines += [
            "\n## 三、决策依据",
        ]
        for basis in decision_basis:
            lines.append(f"- {basis}")

        lines += [
            "\n## 四、图谱核心节点",
        ]
        for req in core_nodes:
            lines.append(f"- [{req.get('priority', '')}] {req.get('name', '')} （得分 {req.get('totalScore', 0)}）")

        lines += [
            "\n## 五、不建议做的需求",
        ]
        for req in not_recommended:
            lines.append(f"- {req.get('name', '')}（得分 {req.get('totalScore', 0)}）：{req.get('aiExplanation', '')}")

        lines += [
            "\n## 六、风险提醒",
        ]
        for warn in risk_warnings:
            lines.append(f"- ⚠️ {warn}")

        filepath.write_text("\n".join(lines), encoding="utf-8")
        logger.info(f"Markdown报告导出完成: {filepath}")
        return str(filepath)

    def export_pdf(self, report_data: dict) -> str:
        """导出 PDF 格式报告，返回文件路径"""
        _ensure_export_dir()
        filename = f"decision_report_{_timestamp()}.pdf"
        filepath = EXPORT_DIR / filename
        logger.info(f"导出PDF报告: {filename}")
    
        try:
            from reportlab.lib.pagesizes import A4
            from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
            from reportlab.lib.units import cm
            from reportlab.lib import colors
            from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
            from reportlab.pdfbase import pdfmetrics
            from reportlab.pdfbase.ttfonts import TTFont
            import platform
            import os
    
            # 尝试注册中文字体（系统字体），按优先级尝试多个字体
            cn_font = "Helvetica"  # 默认 fallback
            font_registered = False
                
            if platform.system() == "Windows":
                # Windows 系统常用中文字体路径列表
                # 格式: (字体文件, subfontIndex) - .ttc文件需要subfontIndex，.ttf文件用None
                cn_font_candidates = [
                    ("C:/Windows/Fonts/msyh.ttc", 0),    # 微软雅黑
                    ("C:/Windows/Fonts/msyhbd.ttc", 0),  # 微软雅黑粗体
                    ("C:/Windows/Fonts/simhei.ttf", None),  # 黑体
                    ("C:/Windows/Fonts/simsun.ttc", 0),     # 宋体
                    ("C:/Windows/Fonts/simkai.ttf", None),  # 楷体
                ]
                for font_path, subfont_idx in cn_font_candidates:
                    if os.path.exists(font_path):
                        try:
                            font_name = "ChineseFont"  # 使用固定名称
                            if subfont_idx is not None:
                                # .ttc 字体集合需要指定 subfontIndex
                                pdfmetrics.registerFont(TTFont(font_name, font_path, subfontIndex=subfont_idx))
                            else:
                                # .ttf 单字体文件
                                pdfmetrics.registerFont(TTFont(font_name, font_path))
                            cn_font = font_name
                            font_registered = True
                            logger.info(f"PDF使用中文字体: {font_path}")
                            break
                        except Exception as e:
                            logger.warning(f"字体注册失败 {font_path}: {e}")
                            continue
                
            if not font_registered:
                logger.warning("未找到中文字体，PDF 可能无法正确显示中文")

            doc = SimpleDocTemplate(str(filepath), pagesize=A4,
                                    leftMargin=2*cm, rightMargin=2*cm,
                                    topMargin=2*cm, bottomMargin=2*cm)
            styles = getSampleStyleSheet()
            # 创建使用中文字体的正文样式，确保所有中文都能正确显示
            normal_style = ParagraphStyle("cn_normal", fontName=cn_font, fontSize=10, spaceAfter=6)
            story = []

            # 标题
            title_style = ParagraphStyle("title", fontName=cn_font, fontSize=18, spaceAfter=12, textColor=colors.HexColor("#1e40af"))
            story.append(Paragraph("需求优先级决策报告", title_style))
            story.append(Paragraph(f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}", normal_style))
            story.append(Spacer(1, 20))

            # 总览
            section_style = ParagraphStyle("section", fontName=cn_font, fontSize=13, spaceBefore=12, spaceAfter=6, textColor=colors.HexColor("#374151"))
            story.append(Paragraph("一、总览", section_style))
            summary = report_data.get("summary", {})
            summary_text = f"需求总数：{summary.get('totalRequirements', 0)}  |  P0需求：{summary.get('p0Count', 0)}  |  高风险：{summary.get('highRiskCount', 0)}  |  平均分：{summary.get('avgScore', 0)}"
            story.append(Paragraph(summary_text, normal_style))
            story.append(Spacer(1, 12))

            # Top10
            story.append(Paragraph("二、本轮优先级 Top 10", section_style))
            top10 = report_data.get("top10", [])
            if top10:
                table_data = [["排名", "需求名称", "优先级", "总分"]]
                for i, req in enumerate(top10, 1):
                    table_data.append([str(i), req.get("name", ""), req.get("priority", ""), str(req.get("totalScore", 0))])
                t = Table(table_data, colWidths=[1.5*cm, 8*cm, 2.5*cm, 2.5*cm])
                t.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor("#dbeafe")),
                    ('FONTNAME', (0, 0), (-1, -1), cn_font),
                    ('FONTSIZE', (0, 0), (-1, -1), 9),
                    ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.HexColor("#f9fafb")]),
                ]))
                story.append(t)
            story.append(Spacer(1, 12))

            # 决策依据
            story.append(Paragraph("三、决策依据", section_style))
            for basis in report_data.get("decisionBasis", []):
                story.append(Paragraph(f"• {basis}", normal_style))
            story.append(Spacer(1, 12))

            # 风险提醒
            story.append(Paragraph("四、风险提醒", section_style))
            for warn in report_data.get("riskWarnings", []):
                story.append(Paragraph(f"⚠ {warn}", normal_style))

            doc.build(story)
            logger.info(f"PDF报告导出完成: {filepath}")
        except ImportError:
            # reportlab 未安装，回退为markdown
            logger.warning("reportlab未安装，回退为Markdown导出")
            return self.export_markdown(report_data)
        except Exception as e:
            logger.error(f"PDF导出失败: {e}")
            return self.export_markdown(report_data)

        return str(filepath)

    def export_pptx(self, report_data: dict) -> str:
        """导出 PPT 格式报告，返回文件路径"""
        _ensure_export_dir()
        filename = f"decision_report_{_timestamp()}.pptx"
        filepath = EXPORT_DIR / filename
        logger.info(f"导出PPTX报告: {filename}")

        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN

            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)

            blank_layout = prs.slide_layouts[6]  # Blank

            def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=(0, 0, 0)):
                txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = text
                run = p.runs[0]
                run.font.size = Pt(font_size)
                run.font.bold = bold
                run.font.color.rgb = RGBColor(*color)
                return txBox

            # Slide 1: 标题页
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "需求优先级决策报告", 1, 2.5, 11, 1.5, font_size=36, bold=True, color=(30, 64, 175))
            summary = report_data.get("summary", {})
            sub = f"需求总数 {summary.get('totalRequirements', 0)} | P0: {summary.get('p0Count', 0)} | 高风险: {summary.get('highRiskCount', 0)} | 平均分: {summary.get('avgScore', 0)}"
            add_text_box(slide, sub, 1, 4.2, 11, 0.6, font_size=14, color=(107, 114, 128))
            add_text_box(slide, f"生成时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}", 1, 5, 11, 0.5, font_size=12, color=(156, 163, 175))

            # Slide 2: Top10
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "本轮优先级 Top 10", 0.5, 0.3, 12, 0.8, font_size=24, bold=True, color=(30, 64, 175))
            top10 = report_data.get("top10", [])
            for i, req in enumerate(top10[:10]):
                row = i // 2
                col = i % 2
                x = 0.5 + col * 6.5
                y = 1.2 + row * 1.1
                text = f"{i+1}. {req.get('name', '')}  [{req.get('priority', '')}]  {req.get('totalScore', 0)}分"
                add_text_box(slide, text, x, y, 6, 0.8, font_size=11)

            # Slide 3: 决策依据 & 风险提醒
            slide = prs.slides.add_slide(blank_layout)
            add_text_box(slide, "决策依据", 0.5, 0.3, 6, 0.6, font_size=20, bold=True, color=(30, 64, 175))
            basis_text = "\n".join(f"• {b}" for b in report_data.get("decisionBasis", []))
            add_text_box(slide, basis_text, 0.5, 1.0, 6, 5.5, font_size=11)

            add_text_box(slide, "风险提醒", 7, 0.3, 6, 0.6, font_size=20, bold=True, color=(220, 38, 38))
            warn_text = "\n".join(f"⚠ {w}" for w in report_data.get("riskWarnings", []))
            add_text_box(slide, warn_text, 7, 1.0, 6, 5.5, font_size=11)

            prs.save(str(filepath))
            logger.info(f"PPTX报告导出完成: {filepath}")
        except ImportError:
            logger.warning("python-pptx未安装，回退为Markdown导出")
            return self.export_markdown(report_data)
        except Exception as e:
            logger.error(f"PPTX导出失败: {e}")
            return self.export_markdown(report_data)

        return str(filepath)


report_exporter = ReportExporter()
